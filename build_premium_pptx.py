import os
import sys
import subprocess

print("[INFO] VERIFYING PYTHON-PPTX INSTALLATION...")
try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

proj_dir = r"C:\Users\rajti\Downloads\Projects\ACADEMIC INTERNSHIP\ClearSight_Project"
pdf_dir = os.path.join(proj_dir, "PDFs")
graphs_dir = os.path.join(pdf_dir, "graphs")
pptx_filepath = os.path.join(pdf_dir, "ClearSight_Premium_Presentation.pptx")

prs = Presentation()
# Widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]

# LIGHT MODE COLORS
BG_COLOR = RGBColor(255, 255, 255) # White
ACCENT_COLOR = RGBColor(15, 60, 150) # Professional Deep Blue
TEXT_COLOR = RGBColor(30, 30, 30) # Dark Charcoal/Black
MUTED_COLOR = RGBColor(80, 80, 80) # Dark Gray

def apply_light_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Arial'

def add_bullet_points(slide, left, top, width, height, points):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Arial'
        p.space_after = Pt(14)
        p.level = 0

# --- SLIDE 1: Title ---
slide_1 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_1)
txBox = slide_1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ClearSight AI"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = ACCENT_COLOR
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Autonomous Kinetic & Biometric Re-Identification Engine"
p2.font.size = Pt(32)
p2.font.color.rgb = TEXT_COLOR
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

txBox2 = slide_1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.3), Inches(1.5))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Submitted By: Rajtilak Chamlagain\nUnder Guidance of: Dr. Mahapara Khursid\nTIH–TIDF, IIT Guwahati"
p3.font.size = Pt(20)
p3.font.color.rgb = MUTED_COLOR
p3.alignment = PP_ALIGN.CENTER

# --- SLIDE 2: The Surveillance Dilemma ---
slide_2 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_2)
add_title(slide_2, "The Surveillance Dilemma")
pts_2 = [
    "Crowd Occlusions: Targets hide behind pillars; basic trackers lose identity.",
    "Environmental Chaos: Shadows and steep angles destroy standard pixel matching.",
    "Static Thresholds: Guessing match limits causes false arrests in the dark.",
    "Browser Overload: Traditional Base64 evidence rendering freezes police laptops."
]
add_bullet_points(slide_2, Inches(1), Inches(2), Inches(11), Inches(4), pts_2)

# --- SLIDE 3: Architecture Pipeline ---
slide_3 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_3)
add_title(slide_3, "5-Phase Architecture Pipeline")
pts_3 = [
    "Phase 1: Biometric Encoding (RetinaFace + ArcFace)",
    "Phase 2: Kinetic Localization (YOLOv8 + ByteTrack)",
    "Phase 3: Unsupervised Selection (Autonomous Spectral Gap Engine)",
    "Phase 4: Evidence Synthesis (3x Slow-Motion Reconstructor)",
    "Phase 5: Zero-Lag UI (Streamlit + Native Disk Sockets)"
]
add_bullet_points(slide_3, Inches(1), Inches(2), Inches(11), Inches(4), pts_3)

# --- SLIDE 4: Solving Crowd Occlusion (With Graph) ---
slide_4 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_4)
add_title(slide_4, "Solving Crowd Occlusion")
pts_4 = [
    "YOLOv8: Real-time bounding box detection.",
    "ByteTrack: Uses Kalman Filters to predict movement during occlusion.",
    "Result: Target IDs stay alive even when completely hidden by crowds!"
]
add_bullet_points(slide_4, Inches(0.5), Inches(2), Inches(5.5), Inches(4), pts_4)
try:
    slide_4.shapes.add_picture(os.path.join(graphs_dir, "Graph_2_Kinetic_Retention.png"), Inches(6.5), Inches(2), width=Inches(6.3))
except:
    pass

# --- SLIDE 5: Biometric Invariance (With Graph) ---
slide_5 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_5)
add_title(slide_5, "Biometric Invariance")
pts_5 = [
    "RetinaFace: 5-point alignment straightens tilted faces automatically.",
    "ArcFace 512D: Rejects FaceNet. Maps faces to a strict mathematical sphere.",
    "Result: 90%+ Accuracy maintained in low light and heavy makeup."
]
add_bullet_points(slide_5, Inches(0.5), Inches(2), Inches(5.5), Inches(4), pts_5)
try:
    slide_5.shapes.add_picture(os.path.join(graphs_dir, "Graph_3_Biometric_Robustness.png"), Inches(6.5), Inches(2), width=Inches(6.3))
except:
    pass

# --- SLIDE 6: Autonomous Spectral Gap (With Graph) ---
slide_6 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_6)
add_title(slide_6, "Autonomous Spectral Gap")
pts_6 = [
    "Humans manually guessing match thresholds is a legal liability.",
    "Algorithm calculates the Maximal Derivative Drop-Off between scores.",
    "Autonomously locks the threshold gate without human bias."
]
add_bullet_points(slide_6, Inches(0.5), Inches(2), Inches(5.5), Inches(4), pts_6)
try:
    slide_6.shapes.add_picture(os.path.join(graphs_dir, "Graph_1_Spectral_Gap.png"), Inches(6.5), Inches(2), width=Inches(6.3))
except:
    pass

# --- SLIDE 7: The Rule of Law ---
slide_7 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_7)
add_title(slide_7, "The Rule of Law (Why NO GenAI?)")
pts_7 = [
    "Generative AI (GFPGAN) hallucinates fake pixels to sharpen blurry faces.",
    "Defense attorneys will successfully argue the evidence is manipulated.",
    "ClearSight strictly computes against authentic, unmanipulated pixels.",
    "Proves we don't need AI enhancements to confidently catch suspects."
]
add_bullet_points(slide_7, Inches(1), Inches(2), Inches(11), Inches(4), pts_7)

# --- SLIDE 8: Zero-Lag Courtroom UI (With Graph) ---
slide_8 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_8)
add_title(slide_8, "Zero-Lag Courtroom UI & Video Triage")
pts_8 = [
    "Native Disk Sockets: Shrinks payload to 0.05MB, stopping browser freezes.",
    "Top-4 Split Strategy: Deep learning models occasionally lose tracking IDs.",
    "We present up to 4 top candidate videos instead of just 1.",
    "Allows human observers to visually verify the suspect from the best 4 angles!"
]
add_bullet_points(slide_8, Inches(0.5), Inches(2), Inches(6.5), Inches(4), pts_8)
try:
    slide_8.shapes.add_picture(os.path.join(graphs_dir, "Graph_4_RAM_Optimization.png"), Inches(7.5), Inches(2), width=Inches(5.0))
except:
    pass


# --- SLIDE 9: Conclusion & Future Scope ---
slide_9 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_9)
add_title(slide_9, "Conclusion & Future Scope")
pts_9 = [
    "Delivered a turnkey digital forensic tool integrating tracking, biometrics, and autonomous math.",
    "Overcame hardware VRAM limits (GTX 1650) via 'Strategic Biometric Stride'.",
    "Future: Scale to Smart-City RTSP networks.",
    "Future: FAISS Vector Databases for nationwide indexing in <15ms."
]
add_bullet_points(slide_9, Inches(1), Inches(2), Inches(11), Inches(4), pts_9)

# --- SLIDE 10: Thank You ---
slide_10 = prs.slides.add_slide(BLANK_LAYOUT)
apply_light_background(slide_10)
tx = slide_10.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
tf = tx.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = ACCENT_COLOR
p.alignment = PP_ALIGN.CENTER

prs.save(pptx_filepath)
print(f"[SUCCESS] LIGHT MODE PPTX SAVED TO: {pptx_filepath}")
