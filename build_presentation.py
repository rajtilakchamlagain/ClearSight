import os
import sys
import subprocess

print("[INFO] VERIFYING PYTHON-PPTX INSTALLATION...")
try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

proj_dir = r"C:\Users\rajti\Downloads\Projects\ACADEMIC INTERNSHIP\ClearSight_Project"
pdf_dir = os.path.join(proj_dir, "PDFs")
os.makedirs(pdf_dir, exist_ok=True)
pptx_filepath = os.path.join(pdf_dir, "ClearSight_AI_Presentation.pptx")

# Create presentation
prs = Presentation()

# Standard layouts
TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
TITLE_AND_CONTENT_LAYOUT = prs.slide_layouts[1]
SECTION_HEADER_LAYOUT = prs.slide_layouts[2]
TWO_CONTENT_LAYOUT = prs.slide_layouts[3]
BLANK_LAYOUT = prs.slide_layouts[6]

# Colors
COLOR_DARK_SLATE = RGBColor(15, 23, 42)
COLOR_BLUE = RGBColor(37, 99, 235)
COLOR_TEXT = RGBColor(51, 65, 85)

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_SLATE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT
    return slide

def add_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(TITLE_AND_CONTENT_LAYOUT)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullets[0]
    
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.color.rgb = COLOR_TEXT
        
    return slide

def add_two_content_slide(prs, title_text, left_bullets, right_bullets):
    slide = prs.slides.add_slide(TWO_CONTENT_LAYOUT)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    left_shape = slide.placeholders[1]
    tf_left = left_shape.text_frame
    tf_left.text = left_bullets[0]
    for bullet in left_bullets[1:]:
        p = tf_left.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.color.rgb = COLOR_TEXT
        
    right_shape = slide.placeholders[2]
    tf_right = right_shape.text_frame
    tf_right.text = right_bullets[0]
    for bullet in right_bullets[1:]:
        p = tf_right.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.color.rgb = COLOR_TEXT
        
    return slide

# --- SLIDE 1: Title Slide ---
subtitle = (
    "A Deep Learning and Unsupervised Thresholding Approach for\n"
    "Court-Admissible Surveillance Person Re-Identification\n\n"
    "Submitted By: Rajtilak Chamlagain\n"
    "Under the Guidance of: Dr. Mahapara Khursid\n"
    "Technology Innovation Hub – TIDF, IIT Guwahati\n"
    "July 2026"
)
add_title_slide(prs, "ClearSight AI: Autonomous Kinetic &\nBiometric Re-Identification Engine", subtitle)

# --- SLIDE 2: Problem Statement ---
bullets_2 = [
    "Modern cities have thousands of CCTV cameras, but finding suspects is manual, slow, and exhausting.",
    "Challenge 1: Heavy Crowd Occlusions. Targets constantly hide behind pillars or other people.",
    "Challenge 2: Lighting & Angle Variations. Standard pixel matching fails under dark or rotated angles.",
    "Challenge 3: The Static Threshold Flaw. Hardcoded matching limits (e.g., 30%) cause wrongful arrests in dim light.",
    "Challenge 4: Browser Overload. Rendering high-res evidence freezes traditional web browsers."
]
add_content_slide(prs, "1. Problem Statement", bullets_2)

# --- SLIDE 3: Project Objectives ---
bullets_3 = [
    "Develop an autonomous deep learning surveillance pipeline for live suspect tracking.",
    "Solve the occlusion problem by adding predictive kinematic memory (ByteTrack) to real-time detection (YOLOv8).",
    "Achieve lighting-invariant facial recognition using 512-Dimensional geometry (ArcFace + RetinaFace).",
    "Eliminate manual human threshold guesswork by inventing an Unsupervised Spectral Gap engine.",
    "Engineer a zero-lag interactive web dashboard capable of courtroom-ready evidential exports."
]
add_content_slide(prs, "2. Project Objectives", bullets_3)

# --- SLIDE 4: Overall Architecture Pipeline ---
left_4 = [
    "Phase 1: Biometric Encoding",
    "  • RetinaFace (5-point alignment)",
    "  • ArcFace (512D Master Signature)",
    "",
    "Phase 2: Kinetic Localization",
    "  • YOLOv8 (Human Bounding Boxes)",
    "  • ByteTrack (Kalman Filter Memory)"
]
right_4 = [
    "Phase 3: Unsupervised Selection",
    "  • Autonomous Spectral Gap Engine",
    "",
    "Phase 4: Evidence Synthesis",
    "  • 3x Slow-Motion Reconstructor",
    "",
    "Phase 5: Zero-Lag Presentation",
    "  • Streamlit + Native Socket Streaming"
]
add_two_content_slide(prs, "3. Overall System Architecture", left_4, right_4)

# --- SLIDE 5: Kinetic Tracking (YOLOv8 + ByteTrack) ---
bullets_5 = [
    "Why YOLOv8?",
    "  • 'You Only Look Once' architecture allows real-time bounding box detection on HD surveillance streams.",
    "Why ByteTrack?",
    "  • Replaces outdated trackers (like DeepSORT) which fail when clothing colors blur.",
    "  • Uses 'Kalman Filter' math to predict a person's trajectory when they walk behind an obstacle.",
    "  • Intelligently recycles low-confidence boxes to keep identities immortal across dense crowds.",
    "[RECOMMENDATION: INSERT A SCREENSHOT/GIF OF YOLO TRACKING HERE]"
]
add_content_slide(prs, "4. Kinetic Localization (YOLOv8 + ByteTrack)", bullets_5)

# --- SLIDE 6: Deep Biometric Engine (RetinaFace + ArcFace) ---
bullets_6 = [
    "RetinaFace (Facial Landmark Alignment):",
    "  • Finds the exact 5 points (eyes, nose, mouth) even on tilted or side-profile CCTV faces.",
    "  • Uses Affine Transformations to mathematically 'straighten' the face.",
    "ArcFace (512D Vector Encoding):",
    "  • Discards outdated 'Triplet Loss' (FaceNet) for modern 'Additive Angular Margin Loss'.",
    "  • Maps faces onto a 512-dimensional mathematical hypersphere.",
    "  • Compares faces via Cosine Similarity (angle between vectors), unaffected by age, lighting, or makeup."
]
add_content_slide(prs, "5. Biometric Encoding (RetinaFace + ArcFace)", bullets_6)

# --- SLIDE 7: The Spectral Gap Engine ---
left_7 = [
    "The Problem with Traditional Models:",
    "  • Humans manually set static thresholds (e.g., 'Match > 50%').",
    "  • In dark scenes, real matches might drop to 30%. In bright scenes, fake matches spike to 60%.",
    "  • Result: Wrongful arrests or missed targets."
]
right_7 = [
    "The Autonomous Spectral Gap Solution:",
    "  • Calculates the largest 'Derivative Drop-Off' between consecutive similarity scores.",
    "  • Automatically places the operational gate inside the massive gap between the suspect and random noise.",
    "  • Zero human intervention. 100% dynamic adaptation."
]
add_two_content_slide(prs, "6. Autonomous Spectral Gap 'Cliff Detection'", left_7, right_7)

# --- SLIDE 8: The Rule of Law (Why No GenAI) ---
bullets_8 = [
    "Why did we reject AI Image Enhancers (GFPGAN) or CLAHE?",
    "  • Generative AI hallucinates and invents fake pixels to make a blurry image look sharp.",
    "  • In a court of law, defense attorneys will immediately invalidate manipulated pixel evidence.",
    "  • ClearSight AI strictly evaluates authentic, unmanipulated pixels using robust deep networks designed for low-resolution conditions.",
    "  • Result: Court-admissible, mathematically verified biometric conclusions."
]
add_content_slide(prs, "7. The Rule of Law (Evidential Integrity)", bullets_8)

# --- SLIDE 9: System Results & Web Application ---
bullets_9 = [
    "Zero-Lag Streamlit Interactive Dashboard:",
    "  • Solved browser JavaScript memory freezes by eliminating 250MB Base64 injections.",
    "  • Implemented native disk socket streaming, achieving 5,000x payload compression.",
    "Widescreen 16:9 Cinema Evidence Galleries:",
    "  • All target snapshots are automatically padded to a 400x225 slate background.",
    "Automated 3x Slow-Motion Synthesis:",
    "  • Generates fractional slow-motion replays for targets appearing < 3.0 seconds."
]
add_content_slide(prs, "8. Web Dashboard & Results", bullets_9)

# --- SLIDE 10: Video Demonstration ---
slide_10 = prs.slides.add_slide(TITLE_AND_CONTENT_LAYOUT)
title_10 = slide_10.shapes.title
title_10.text = "9. Live Demonstration / Video Proof"
title_10.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
title_10.text_frame.paragraphs[0].font.bold = True
tf_10 = slide_10.placeholders[1].text_frame
tf_10.text = "INSERT YOUR FINAL SURVEILLANCE MP4 VIDEO HERE."
p10 = tf_10.add_paragraph()
p10.text = "(Go to Insert -> Video -> This Device in PowerPoint)"
p10.level = 0

# --- SLIDE 11: Challenges Faced ---
bullets_11 = [
    "Hardware VRAM Constraints:",
    "  • Limited to GTX 1650 (4GB VRAM).",
    "  • Solved via 'Strategic Biometric Stride' (Running ArcFace only every 3rd frame).",
    "Windows Certificate Store (SSL) Errors:",
    "  • Web server crashed due to Windows SSL proxy exceptions.",
    "  • Engineered a guarded 'start.py' wrapper to intercept SSL certificates before Streamlit load.",
    "Dataset Limitations:",
    "  • Sourced unconstrained raw political rallies and VIP walkthroughs to simulate real police operations."
]
add_content_slide(prs, "10. Challenges Overcome", bullets_11)

# --- SLIDE 12: Conclusion & Future Scope ---
bullets_12 = [
    "Conclusion:",
    "  • Engineered a complete, turnkey digital forensic surveillance platform.",
    "  • Eliminated threshold guesswork and occlusion-tracking failures.",
    "Future Scope:",
    "  • Multi-Camera RTSP Integration: Expand from video files to live smart-city network streams.",
    "  • FAISS Vector Databases: Enable nationwide instantaneous identity searches in < 15 milliseconds.",
    "  • 3D Gait Reconstruction: Add volumetric pose regression for enhanced courtroom presentation."
]
add_content_slide(prs, "11. Conclusion & Future Scope", bullets_12)

# --- SLIDE 13: Q&A ---
slide_13 = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
title_13 = slide_13.shapes.title
title_13.text = "Thank You"
title_13.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_SLATE
title_13.text_frame.paragraphs[0].font.bold = True
subtitle_13 = slide_13.placeholders[1]
subtitle_13.text = "Questions & Answers"

prs.save(pptx_filepath)
print(f"[SUCCESS] PRESENTATION SAVED TO: {pptx_filepath}")
